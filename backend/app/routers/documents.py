"""Document ingestion — parse, classify, enhance, and memorise.

Supported formats:
  Documents : PDF, DOCX, PPTX, TXT, MD, CSV
  Images    : PNG, JPG, JPEG, WEBP, BMP, GIF  ← text extracted via Groq vision
  Video/Audio: MP4, MOV, WEBM, MP3, WAV, M4A  ← metadata only (Whisper opt-in)

On every upload the backend:
  1. Extracts text (parser per mime-type, Groq vision for images)
  2. Runs LLM classification → doc_type, subject, summary, key_topics
  3. For images: Pillow auto-enhance + rembg background removal
  4. Chunks text into MemoryRecord rows with role="document"

Extra endpoints:
  POST /documents/{id}/enhance    → re-run Pillow enhance on an image
  POST /documents/{id}/remove-bg  → re-run rembg on an image
"""

from __future__ import annotations

import base64
import io
import json
import logging
import mimetypes
import re
from datetime import datetime
from typing import Optional

from cuid2 import Cuid as _Cuid
_cuid_gen = _Cuid()
def cuid() -> str:
    return _cuid_gen.generate()

from fastapi import APIRouter, Depends, File, HTTPException, UploadFile
from pydantic import BaseModel
from sqlmodel import Session, select

from ..agent.agent_loop import get_async_client
from ..config import get_settings
from ..db import get_session
from ..deps import get_current_user
from ..models.memory import MemoryRecord
from ..models.user import User
from ..services.semantic_memory import _heuristic_importance as compute_importance

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/documents", tags=["documents"])

# ── Limits ────────────────────────────────────────────────────────────────────
MAX_BYTES   = 50 * 1024 * 1024
CHUNK_SIZE  = 1200
CHUNK_OVERLAP = 200

# ── MIME sets ─────────────────────────────────────────────────────────────────
_PDF_TYPES   = {"application/pdf"}
_DOCX_TYPES  = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
}
_PPTX_TYPES  = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}
_TEXT_TYPES  = {"text/plain", "text/markdown", "text/csv", "application/json"}
_IMAGE_TYPES = {
    "image/png", "image/jpeg", "image/jpg", "image/webp",
    "image/bmp", "image/gif", "image/tiff",
}
_VIDEO_TYPES = {"video/mp4", "video/mpeg", "video/webm", "video/ogg", "video/quicktime"}
_AUDIO_TYPES = {"audio/mpeg", "audio/wav", "audio/ogg", "audio/mp4", "audio/webm"}

ALL_ACCEPTED = (
    _PDF_TYPES | _DOCX_TYPES | _PPTX_TYPES | _TEXT_TYPES
    | _IMAGE_TYPES | _VIDEO_TYPES | _AUDIO_TYPES
)

# ── Text parsers ──────────────────────────────────────────────────────────────

def _parse_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join(
        (page.extract_text() or "").strip() for page in reader.pages
    ).strip()


def _parse_docx(data: bytes) -> str:
    from docx import Document  # type: ignore
    doc = Document(io.BytesIO(data))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation  # type: ignore
    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            " ".join(run.text for run in para.runs).strip()
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs
        ]
        texts = [t for t in texts if t]
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _parse_text(data: bytes) -> str:
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


async def _parse_image(data: bytes, mime: str, filename: str) -> str:
    """Use Groq vision to describe the image and extract any visible text."""
    try:
        b64 = base64.b64encode(data).decode()
        data_url = f"data:{mime};base64,{b64}"
        settings = get_settings()
        client = get_async_client()
        completion = await client.chat.completions.create(
            model=settings.fast_llm_model or settings.llm_model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {"url": data_url},
                        },
                        {
                            "type": "text",
                            "text": (
                                "Please do two things:\n"
                                "1. Extract every piece of text visible in this image verbatim.\n"
                                "2. Describe what the image shows (diagram, chart, photo, slide, etc.).\n"
                                "Format:\n"
                                "EXTRACTED TEXT:\n<text or 'None'>\n\n"
                                "DESCRIPTION:\n<description>"
                            ),
                        },
                    ],
                }
            ],
            max_tokens=1024,
            temperature=0.1,
        )
        return (completion.choices[0].message.content or "").strip()
    except Exception as exc:
        logger.warning("Vision OCR failed for %s: %s", filename, exc)
        return f"[Image file: {filename}. Vision OCR unavailable: {exc}]"


def _parse_video_audio(filename: str, size_bytes: int) -> str:
    return (
        f"[Media file: {filename}]\n"
        f"Size: {size_bytes / 1024 / 1024:.1f} MB\n"
        "Full speech-to-text transcription requires a Whisper API key (WHISPER_API_KEY)."
    )


async def _extract_text(data: bytes, mime: str, filename: str) -> str:
    try:
        if mime in _PDF_TYPES:   return _parse_pdf(data)
        if mime in _DOCX_TYPES:  return _parse_docx(data)
        if mime in _PPTX_TYPES:  return _parse_pptx(data)
        if mime in _TEXT_TYPES:  return _parse_text(data)
        if mime in _IMAGE_TYPES: return await _parse_image(data, mime, filename)
        if mime in (_VIDEO_TYPES | _AUDIO_TYPES):
            return _parse_video_audio(filename, len(data))
    except Exception as exc:
        logger.warning("Parser error for %s (%s): %s", filename, mime, exc)
        return f"[Parse failed for {filename}: {exc}]"
    return f"[Unsupported: {mime}]"


# ── Image processing helpers ──────────────────────────────────────────────────

def _pillow_enhance(data: bytes) -> bytes:
    """Auto-enhance: contrast + sharpness + subtle brightness boost."""
    from PIL import Image, ImageEnhance, ImageFilter  # type: ignore
    img = Image.open(io.BytesIO(data)).convert("RGBA")

    # Work on RGB layer only, preserve alpha
    rgb = img.convert("RGB")
    rgb = ImageEnhance.Contrast(rgb).enhance(1.3)
    rgb = ImageEnhance.Sharpness(rgb).enhance(1.4)
    rgb = ImageEnhance.Color(rgb).enhance(1.15)
    rgb = ImageEnhance.Brightness(rgb).enhance(1.05)
    # Mild unsharp-mask for crispness
    rgb = rgb.filter(ImageFilter.UnsharpMask(radius=1.5, percent=80, threshold=3))

    # Merge back alpha if present
    r, g, b = rgb.split()
    _, _, _, a = img.split()
    out = Image.merge("RGBA", (r, g, b, a))

    buf = io.BytesIO()
    out.save(buf, format="PNG")
    return buf.getvalue()


def _remove_background(data: bytes) -> bytes:
    """Remove background using rembg (U2Net). Returns RGBA PNG."""
    from rembg import remove as rembg_remove  # type: ignore
    result: bytes = rembg_remove(data)
    return result


def _thumbnail_b64(data: bytes, mime: str, max_px: int = 300) -> str:
    """Resize image to thumbnail and return as base64 PNG data-URL."""
    from PIL import Image  # type: ignore
    img = Image.open(io.BytesIO(data))
    img.thumbnail((max_px, max_px), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    b64 = base64.b64encode(buf.getvalue()).decode()
    return f"data:image/png;base64,{b64}"


# ── LLM classification ────────────────────────────────────────────────────────

_CLASSIFY_PROMPT = """You are a document classifier. Given the first ~1500 characters of a document,
return STRICT JSON with exactly this shape:

{{
  "doc_type": "<one of: Research Paper, Lecture Notes, Textbook Chapter, Presentation, Study Guide, Code/Technical, Image/Diagram, Video Transcript, General Notes, Other>",
  "subject": "<subject area, e.g. Machine Learning, Calculus, History, Chemistry>",
  "summary": "<2-3 sentence plain-English summary of what this document covers>",
  "key_topics": ["<topic>", "<topic>", "<topic>"]
}}

Document excerpt:
{excerpt}
"""


async def _classify_document(text: str) -> dict:
    try:
        settings = get_settings()
        client = get_async_client()
        excerpt = text[:1500].strip()
        if not excerpt:
            raise ValueError("Empty text")
        completion = await client.chat.completions.create(
            model=settings.fast_llm_model or settings.llm_model,
            messages=[
                {"role": "system", "content": "You output strict JSON only. No markdown, no code fences."},
                {"role": "user", "content": _CLASSIFY_PROMPT.format(excerpt=excerpt)},
            ],
            max_tokens=400,
            temperature=0.1,
        )
        raw = (completion.choices[0].message.content or "").strip()
        # Strip fences if any
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        return json.loads(raw)
    except Exception as exc:
        logger.warning("Classification failed: %s", exc)
        return {
            "doc_type": "Other",
            "subject": "Unknown",
            "summary": "",
            "key_topics": [],
        }


# ── Chunker ───────────────────────────────────────────────────────────────────

def _chunk(text: str) -> list[str]:
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not text:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + CHUNK_SIZE, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = end - CHUNK_OVERLAP
    return chunks


# ── Response schemas ──────────────────────────────────────────────────────────

class DocumentOut(BaseModel):
    id: str
    filename: str
    mime: str
    size_bytes: int
    chunk_count: int
    preview: str
    created_at: datetime
    # Classification
    doc_type: str = "Other"
    subject: str = ""
    summary: str = ""
    key_topics: list[str] = []
    # Image extras
    is_image: bool = False
    thumbnail_b64: Optional[str] = None
    enhanced_b64: Optional[str] = None
    nobg_b64: Optional[str] = None


class DocumentDetail(DocumentOut):
    chunks: list[str]


def _meta_to_out(record: MemoryRecord) -> Optional[DocumentOut]:
    try:
        meta = json.loads(record.content)
        if not meta.get("_doc_meta"):
            return None
        return DocumentOut(
            id=record.id,
            filename=meta.get("filename", "unknown"),
            mime=meta.get("mime", ""),
            size_bytes=meta.get("size_bytes", 0),
            chunk_count=meta.get("chunk_count", 0),
            preview=meta.get("preview", ""),
            created_at=record.created_at,
            doc_type=meta.get("doc_type", "Other"),
            subject=meta.get("subject", ""),
            summary=meta.get("summary", ""),
            key_topics=meta.get("key_topics", []),
            is_image=meta.get("is_image", False),
            thumbnail_b64=meta.get("thumbnail_b64"),
            enhanced_b64=meta.get("enhanced_b64"),
            nobg_b64=meta.get("nobg_b64"),
        )
    except Exception:
        return None


# ── Upload endpoint ───────────────────────────────────────────────────────────

@router.post("/upload", response_model=DocumentOut)
async def upload_document(
    file: UploadFile = File(...),
    session: Session = Depends(get_session),
    user: User = Depends(get_current_user),
) -> DocumentOut:
    data = await file.read()
    if len(data) > MAX_BYTES:
        raise HTTPException(status_code=413, detail=f"File too large (max {MAX_BYTES // (1024*1024)} MB)")

    filename = file.filename or "upload"
    mime = file.content_type or mimetypes.guess_type(filename)[0] or "application/octet-stream"

    # Normalise common aliases
    ext_lower = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    ext_map = {
        "md": "text/markdown", "csv": "text/csv", "txt": "text/plain",
        "png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
        "webp": "image/webp", "bmp": "image/bmp", "gif": "image/gif",
        "heic": "image/heic", "tiff": "image/tiff",
    }
    if ext_lower in ext_map:
        mime = ext_map[ext_lower]

    if mime not in ALL_ACCEPTED:
        raise HTTPException(
            status_code=415,
            detail=f"Unsupported file type '{mime}'. Accepted: PDF, DOCX, PPTX, TXT, MD, CSV, PNG, JPG, WEBP, MP4, MP3, WAV…",
        )

    is_image = mime in _IMAGE_TYPES

    # ── 1. Extract text ───────────────────────────────────────────────────────
    text = await _extract_text(data, mime, filename)
    chunks = _chunk(text)
    if not chunks:
        chunks = [f"[Empty document: {filename}]"]

    # ── 2. Classify ───────────────────────────────────────────────────────────
    classification = await _classify_document(text)

    # ── 3. Image processing ───────────────────────────────────────────────────
    thumbnail_b64: Optional[str] = None
    enhanced_b64: Optional[str] = None
    nobg_b64: Optional[str] = None

    if is_image:
        try:
            thumbnail_b64 = _thumbnail_b64(data, mime)
        except Exception as exc:
            logger.warning("Thumbnail failed: %s", exc)

        try:
            enhanced_bytes = _pillow_enhance(data)
            enhanced_b64 = f"data:image/png;base64,{base64.b64encode(enhanced_bytes).decode()}"
        except Exception as exc:
            logger.warning("Enhance failed: %s", exc)

        try:
            nobg_bytes = _remove_background(data)
            nobg_b64 = f"data:image/png;base64,{base64.b64encode(nobg_bytes).decode()}"
        except Exception as exc:
            logger.warning("BG removal failed: %s", exc)

    # ── 4. Store in DB ────────────────────────────────────────────────────────
    doc_id = cuid()
    preview = text[:300].replace("\n", " ").strip()

    meta_content = json.dumps({
        "_doc_meta": True,
        "filename": filename,
        "mime": mime,
        "size_bytes": len(data),
        "chunk_count": len(chunks),
        "preview": preview,
        # Classification
        "doc_type": classification.get("doc_type", "Other"),
        "subject": classification.get("subject", ""),
        "summary": classification.get("summary", ""),
        "key_topics": classification.get("key_topics", []),
        # Image extras (full data-URLs stored for instant retrieval)
        "is_image": is_image,
        "thumbnail_b64": thumbnail_b64,
        "enhanced_b64": enhanced_b64,
        "nobg_b64": nobg_b64,
    })

    meta_record = MemoryRecord(
        id=doc_id,
        user_id=user.id,
        role="document",
        session_id=doc_id,
        content=meta_content,
        importance=0.9,
    )
    session.add(meta_record)

    for i, chunk in enumerate(chunks):
        session.add(MemoryRecord(
            id=cuid(),
            user_id=user.id,
            role="document",
            session_id=doc_id,
            content=f"[From '{filename}', part {i+1}/{len(chunks)}]\n{chunk}",
            importance=compute_importance(chunk),
        ))

    session.commit()
    session.refresh(meta_record)

    return DocumentOut(
        id=doc_id,
        filename=filename,
        mime=mime,
        size_bytes=len(data),
        chunk_count=len(chunks),
        preview=preview,
        created_at=meta_record.created_at,
        doc_type=classification.get("doc_type", "Other"),
        subject=classification.get("subject", ""),
        summary=classification.get("summary", ""),
        key_topics=classification.get("key_topics", []),
        is_image=is_image,
        thumbnail_b64=thumbnail_b64,
        enhanced_b64=enhanced_b64,
        nobg_b64=nobg_b64,
    )


# ── Image action endpoints ────────────────────────────────────────────────────

class ImageActionResponse(BaseModel):
    result_b64: str   # data:image/png;base64,…


@router.post("/{doc_id}/enhance", response_model=ImageActionResponse)
async def enhance_image(
    doc_id: str,
    session: Session = Depends(get_session),
    user: User = Depends(get_current_user),
) -> ImageActionResponse:
    meta_record = session.get(MemoryRecord, doc_id)
    if not meta_record or meta_record.user_id != user.id:
        raise HTTPException(status_code=404, detail="Document not found")
    try:
        meta = json.loads(meta_record.content)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid document metadata")

    # If we already have enhanced_b64, return it directly
    if meta.get("enhanced_b64"):
        return ImageActionResponse(result_b64=meta["enhanced_b64"])

    raise HTTPException(status_code=400, detail="Not an image document or enhancement unavailable")


@router.post("/{doc_id}/remove-bg", response_model=ImageActionResponse)
async def remove_bg(
    doc_id: str,
    session: Session = Depends(get_session),
    user: User = Depends(get_current_user),
) -> ImageActionResponse:
    meta_record = session.get(MemoryRecord, doc_id)
    if not meta_record or meta_record.user_id != user.id:
        raise HTTPException(status_code=404, detail="Document not found")
    try:
        meta = json.loads(meta_record.content)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid document metadata")

    if meta.get("nobg_b64"):
        return ImageActionResponse(result_b64=meta["nobg_b64"])

    raise HTTPException(status_code=400, detail="Not an image document or background removal unavailable")


# ── List / get / delete ───────────────────────────────────────────────────────

@router.get("", response_model=list[DocumentOut])
def list_documents(
    session: Session = Depends(get_session),
    user: User = Depends(get_current_user),
) -> list[DocumentOut]:
    stmt = (
        select(MemoryRecord)
        .where(MemoryRecord.user_id == user.id)
        .where(MemoryRecord.role == "document")
        .order_by(MemoryRecord.created_at.desc())
    )
    rows = list(session.exec(stmt))
    docs: list[DocumentOut] = []
    seen: set[str] = set()
    for r in rows:
        if r.session_id and r.session_id not in seen:
            out = _meta_to_out(r)
            if out:
                seen.add(r.session_id)
                docs.append(out)
    return docs


@router.get("/{doc_id}", response_model=DocumentDetail)
def get_document(
    doc_id: str,
    session: Session = Depends(get_session),
    user: User = Depends(get_current_user),
) -> DocumentDetail:
    meta_record = session.get(MemoryRecord, doc_id)
    if not meta_record or meta_record.user_id != user.id:
        raise HTTPException(status_code=404, detail="Document not found")
    out = _meta_to_out(meta_record)
    if not out:
        raise HTTPException(status_code=404, detail="Document not found")
    chunk_stmt = (
        select(MemoryRecord)
        .where(MemoryRecord.user_id == user.id)
        .where(MemoryRecord.session_id == doc_id)
        .where(MemoryRecord.id != doc_id)
        .order_by(MemoryRecord.created_at)
    )
    chunks = [r.content for r in session.exec(chunk_stmt)]
    return DocumentDetail(**out.model_dump(), chunks=chunks)


@router.delete("/{doc_id}")
def delete_document(
    doc_id: str,
    session: Session = Depends(get_session),
    user: User = Depends(get_current_user),
) -> dict:
    stmt = (
        select(MemoryRecord)
        .where(MemoryRecord.user_id == user.id)
        .where(MemoryRecord.session_id == doc_id)
    )
    rows = list(session.exec(stmt))
    if not rows:
        raise HTTPException(status_code=404, detail="Document not found")
    for r in rows:
        session.delete(r)
    session.commit()
    return {"ok": True, "deleted": len(rows)}
